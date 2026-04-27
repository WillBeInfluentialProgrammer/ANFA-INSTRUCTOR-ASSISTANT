from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define ANF Colors
ANF_GREEN = RGBColor(20, 83, 45)
ANF_LIGHT_GREEN = RGBColor(246, 255, 248)
ANF_ACCENT = RGBColor(30, 124, 76)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(50, 50, 50)

def add_title_slide(prs, title_text):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = ANF_GREEN
    
    # Add title
    left = Inches(0.5)
    top = Inches(3)
    width = Inches(9)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    p = title_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title_text, content_list):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.4)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ANF_GREEN
    
    # Add green line under title
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.3), Inches(9.5), Inches(1.3))
    line.line.color.rgb = ANF_ACCENT
    line.line.width = Pt(3)
    
    # Add content
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(8)
    height = Inches(5)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)
        p.space_after = Pt(12)
        p.level = 0
    
    return slide

# =====================================================
# SLIDE 1: TITLE SLIDE
# =====================================================
add_title_slide(prs, "Instructor Assistant")

# =====================================================
# SLIDE 2: PROBLEM STATEMENT
# =====================================================
problem_content = [
    "❌ Manual question creation is time-consuming and labor-intensive",
    "❌ Inconsistent quality in assessments across trainers",
    "❌ Limited context-awareness for legal provisions and sections",
    "❌ No intelligent retrieval system for specific legal references",
    "❌ Risk of AI hallucinations generating inaccurate information",
    "❌ Difficulty in processing and managing multiple PDF documents"
]
add_content_slide(prs, "Problem Statement", problem_content)

# =====================================================
# SLIDE 3: PROPOSED SOLUTION
# =====================================================
solution_content = [
    "✅ AI-Powered RAG (Retrieval Augmented Generation) Architecture",
    "✅ Hybrid Search: Semantic + BM25 keyword matching",
    "✅ Multi-format content generation: MCQs, Descriptive Questions, Scenarios",
    "✅ Streamlit UI for easy instructor interaction",
    "✅ LangChain + Groq LLaMA 3.1 for intelligent responses",
    "✅ Chroma Vector DB with HuggingFace embeddings for context grounding"
]
add_content_slide(prs, "Proposed Solution", solution_content)

# =====================================================
# SLIDE 4: IMPACT & APPLICATIONS
# =====================================================
impact_content = [
    "📚 For Instructors: 80% reduction in question creation time",
    "🎓 For Learners: Realistic scenario-based learning experiences",
    "⚖️ For Legal Training: 100% source-grounded responses (no hallucinations)",
    "🚀 For Organizations: Scalable, cost-effective training solution",
    "🔍 Smart Document Matching: Automatic section extraction and retrieval",
    "🛡️ Applications: ANF training, legal education, regulatory compliance training"
]
add_content_slide(prs, "Impact & Applications", impact_content)

# =====================================================
# SLIDE 5: CONCLUSION & OUTCOMES
# =====================================================
conclusion_content = [
    "✨ Fully functional AI-powered teaching assistant system",
    "📊 100% source attribution - all answers backed by documents",
    "🎯 Three content generation modes: MCQ, Descriptive, Scenario-based",
    "💻 Production-ready with Streamlit web interface",
    "🔄 Future Enhancement: Multi-language support, advanced analytics",
    "🏆 Result: Transforming legal and regulatory training delivery"
]
add_content_slide(prs, "Conclusion & Outcomes", conclusion_content)

# Save presentation
prs.save('ANFA_Instructor_Assistant_Presentation.pptx')
print("✅ Presentation created successfully!")
print("📁 File: ANFA_Instructor_Assistant_Presentation.pptx")
